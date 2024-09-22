#!/usr/bin/env python3
#
# The Week Ahead Interpretation in the Songjiang Daily Bulletin Build System
# Copyright (C) 2024 Runxi Yu <https://runxiyu.org>
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU Affero General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU Affero General Public License for more details.
#
# You should have received a copy of the GNU Affero General Public License
# along with this program.  If not, see <https://www.gnu.org/licenses/>.
#

import logging
import datetime
import os

import pptx

import common

logger = logging.getLogger(__name__)

def download_or_report_the_week_ahead(token: str, datetime_target: datetime.datetime, the_week_ahead_url: str) -> None:
    the_week_ahead_filename = "the_week_ahead-%s.pptx" % datetime_target.strftime("%Y%m%d")
    if not os.path.isfile(the_week_ahead_filename):
        logger.info("Downloading The Week Ahead to %s" % the_week_ahead_filename)
        common.download_share_url(token, the_week_ahead_url, the_week_ahead_filename)
        assert os.path.isfile(the_week_ahead_filename)
    else:
        logger.info("The Week Ahead already exists at %s" % the_week_ahead_filename)

def parse_the_week_ahead(datetime_target: datetime.datetime, the_week_ahead_community_time_page_number: int, the_week_ahead_aod_page_number: int) -> tuple[list[list[str]], list[str]]:
    logger.info("Parsing The Week Ahead")
    the_week_ahead_filename = "the_week_ahead-%s.pptx" % datetime_target.strftime("%Y%m%d")
    the_week_ahead_presentation = pptx.Presentation(the_week_ahead_filename)
    community_time = extract_community_time(
        the_week_ahead_presentation,
        the_week_ahead_community_time_page_number,
    )
    aods = extract_aods(the_week_ahead_presentation, the_week_ahead_aod_page_number)
    return community_time, aods

def extract_community_time(prs: pptx.presentation.Presentation, community_time_page_number: int) -> list[list[str]]:
    slide = prs.slides[community_time_page_number]
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        break
    else:
        raise ValueError("No shapes")
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    if col_count not in [4, 5]:
        raise ValueError("Community time parsing: The Week Ahead community time table does not have 4 or 5 columns")
    if col_count == 4:
        logger.warning("Community time warning: only four columns found, assuming that Y12 has graduated")

    res = [["" for c in range(col_count)] for r in range(row_count)]

    for r in range(row_count):
        for c in range(col_count):
            cell = tbl.cell(r, c)
            if not cell.is_spanned:
                t = ""
                for p in cell.text_frame.paragraphs:
                    for pr in p.runs:
                        t += pr.text
                t = t.strip()
                if "whole school assembly" in t.lower():
                    t = "Whole School Assembly"
                elif "tutor group check-in" in t.lower() or "follow up day" in t.lower() or "open session for tutor and tutee" in t.lower():
                    t = "Tutor Time"
                res[r][c] = t
                if cell.is_merge_origin:
                    for sh in range(cell.span_height):
                        for sw in range(cell.span_width):
                            res[r + sh][c + sw] = t

    return [x[1:] for x in res[1:]]

def extract_aods(prs: pptx.presentation.Presentation, aod_page_number: int) -> list[str]:
    slide = prs.slides[aod_page_number]
    aods = ["", "", "", ""]
    for shape in slide.shapes:
        if hasattr(shape, "text") and "monday: " in shape.text.lower():
            slist = shape.text.split("\n")
            for s in slist:
                try:
                    day, aod = s.split(": ", 1)
                except ValueError:
                    pass
                day = day.lower()
                if day == "monday":
                    aods[0] = aod
                elif day == "tuesday":
                    aods[1] = aod
                elif day == "wednesday":
                    aods[2] = aod
                elif day == "thursday":
                    aods[3] = aod
            if not all(aods):
                raise common.DailyBulletinError("The Week Ahead doesn't include all AOD days, or the formatting is borked")
            return aods
    raise common.DailyBulletinError("The Week Ahead's doesn't even include an AOD for Monday")
    # TODO: this is one of those places where Monday is *expected* to be the first day.
    # TODO: revamp this. this is ugly!
