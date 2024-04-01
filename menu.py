from __future__ import annotations
from typing import Optional, Any
import os
from configparser import ConfigParser
from datetime import datetime
import json
import tempfile
import pptx  # type: ignore
import pptx.exc
from pprint import pprint

DAYNAMES = ["Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"]


def parse_slide(slide: pptx.slide):
    for shape in slide.shapes:
        if shape.has_table:
            break
    else:
        raise ValueError("Slide doesn't contain any tables?")
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    tbll = []
    for r in range(row_count):
        row = [""] * col_count
        old_cell_text = ""
        for c in range(col_count):
            cell_text = ""
            cell = tbl.cell(r, c)
            paragraphs = cell.text_frame.paragraphs
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    cell_text += run.text
            row[c] = (
                "o" if cell.is_merge_origin else ("s" if cell.is_spanned else "n"),
                cell.span_height,
                cell.span_width,
                cell_text.strip(),
            )
            old_cell_text = cell_text
        tbll.append(row)
    return tbll


def zero_list(l):
    return [
        (
            zero_list(i)
            if (type(i) is list or type(i) is tuple or type(i) is set)
            else ""
        )
        for i in l
    ]


def equal_shapes(a, b):
    return zero_list(a) == zero_list(b)


def parse_meal_tables(tbl):
    windows = []
    for j in range(1, len(tbl)):
        cell = tbl[j][0]
        if cell[0] in ["o", "n"]:
            windows.append((j, j - 1 + cell[1]))

    days = ["Mon", "Tue", "Wed", "Thu", "Fri"]
    daysmenus = [[], [], [], [], []]

    assert len(tbl[0]) == 6

    for i in range(1, len(tbl[0])):
        for s, f in windows:
            thiswindow = []
            for j in range(s, f + 1):
                if (
                    tbl[j][i][-1].strip()
                    and tbl[j][i][-1].strip().lower() != "condiments selection"
                ):
                    thiswindow.append(tbl[j][i][-1])
            daysmenus[i - 1].append(thiswindow)
    return daysmenus


def combine_parsed_meal_tables(en, cn):
    if not equal_shapes(cn, en):
        raise ValueError("Augmented menus not in the same shape")

    c = zero_list(en)

    for j in range(len(en)):
        for i in range(len(en[j])):
            for k in range(len(en[j][i])):
                c[j][i][k] = en[j][i][k] + "\n" + cn[j][i][k]
    return c


def extract_menu(
    filename_en: str, filename_cn: str, config: ConfigParser
) -> list[list[str]]:
    try:
        enprs = pptx.Presentation(filename_en)
        cnprs = pptx.Presentation(filename_cn)
    except pptx.exc.PackageNotFoundError:
        raise ValueError(
            "Presentation path %s doesn't exist or is broken" % filename
        ) from None

    breakfast = combine_parsed_meal_tables(
        parse_meal_tables(
            parse_slide(
                enprs.slides[int(config["weekly_menu"]["breakfast_page_number"])]
            )
        ),
        parse_meal_tables(
            parse_slide(
                cnprs.slides[int(config["weekly_menu"]["breakfast_page_number"])]
            )
        ),
    )
    lunch = combine_parsed_meal_tables(
        parse_meal_tables(
            parse_slide(enprs.slides[int(config["weekly_menu"]["lunch_page_number"])])
        ),
        parse_meal_tables(
            parse_slide(cnprs.slides[int(config["weekly_menu"]["lunch_page_number"])])
        ),
    )
    dinner = combine_parsed_meal_tables(
        parse_meal_tables(
            parse_slide(enprs.slides[int(config["weekly_menu"]["dinner_page_number"])])
        ),
        parse_meal_tables(
            parse_slide(cnprs.slides[int(config["weekly_menu"]["dinner_page_number"])])
        ),
    )
    return breakfast, lunch, dinner


def main() -> None:
    today = datetime.today().strftime("%Y%m%d")
    config = ConfigParser()
    config.read("config.ini")
    open("build/menu.json", "w").write(
        json.dumps(
            extract_menu(
                "build/20240408-menu-en.pptx", "build/20240408-menu-cn.pptx", config
            ),
            ensure_ascii=False,
            indent=4,
        )
    )


if __name__ == "__main__":
    main()
