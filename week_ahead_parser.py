from __future__ import annotations
from pprint import pprint  # FIXME
import logging
import msal  # type: ignore
import requests
import datetime
import os
from configparser import ConfigParser
from typing import Any, Optional
import json
import base64
from pptx import Presentation  # type: ignore
import pptx.exc  # type: ignore


DAYNAMES = ["Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"]


def acquire_token(app: msal.PublicClientApplication, config: ConfigParser) -> str:
    result = app.acquire_token_by_username_password(
        config["credentials"]["username"],
        config["credentials"]["password"],
        scopes=config["credentials"]["scope"].split(" "),
    )

    if "access_token" in result:
        assert type(result["access_token"]) is str
        return result["access_token"]
    else:
        raise ValueError("Authentication error in password login")


def encode_sharing_url(url: str) -> str:
    return "u!" + base64.urlsafe_b64encode(url.encode("utf-8")).decode("ascii").rstrip(
        "="
    )


def download_share_url(token: str, url: str, local_filename: str) -> None:
    download_direct_url = requests.get(
        "https://graph.microsoft.com/v1.0/shares/%s/driveItem"
        % encode_sharing_url(url),
        headers={"Authorization": "Bearer " + token},
    ).json()["@microsoft.graph.downloadUrl"]
    r = requests.get(download_direct_url, stream=True)
    with open(local_filename, "wb") as f:
        for chunk in r.iter_content(chunk_size=1024**2):
            if chunk:
                print(".")
                f.write(chunk)


def extract_community_time_from_presentation(config: ConfigParser) -> list[list[str]]:

    try:
        prs = Presentation(config["the_week_ahead"]["local_filename"])
    except pptx.exc.PackageNotFoundError:
        raise ValueError("The Week Ahead is missing, empty, or broken") from None

    slide = prs.slides[int(config["the_week_ahead"]["community_time_page_number"])]
    for shape in slide.shapes:
        if not shape.has_table:
            continue
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    if row_count != 5 or col_count != 5:
        raise ValueError(
            "Community time parsing: The Week Ahead community time table is not 5x5"
        )
    tbll = []
    for r in range(row_count):
        row = [""] * col_count
        for c in range(col_count):
            cell_text = ""
            cell = tbl.cell(r, c)
            paragraphs = cell.text_frame.paragraphs
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    cell_text += run.text
            if not cell_text.strip():
                cell_text = old_cell_text  # type: ignore
                # TODO: Use cell.is_merge_origin, cell.is_spanned,
                # cell.span_height, cell.span_width instead
            row[c] = cell_text
            old_cell_text = cell_text
        tbll.append(row)
    return tbll


def extract_aod_from_presentation(config: ConfigParser) -> list[str]:

    prs = Presentation(config["the_week_ahead"]["local_filename"])
    slide = prs.slides[int(config["the_week_ahead"]["aod_page_number"])]

    aods = ["", "", "", ""]
    for shape in slide.shapes:
        if hasattr(shape, "text") and "Monday: " in shape.text:
            slist = shape.text.split("\n")
            for s in slist:
                try:
                    day, aod = s.split(": ", 1)
                except ValueError:
                    pass
                day = day.lower()
                if day == "monday":
                    aods[0] = aod
                elif day == "tuesday":
                    aods[1] = aod
                elif day == "wednesday":
                    aods[2] = aod
                elif day == "thursday":
                    aods[3] = aod
            if not all(aods):
                raise ValueError(
                    "AOD parsing: The Week Ahead doesn't include all AOD days, or the formatting is borked"
                )
            return aods
            break
    else:
        raise ValueError(
            "AOD parsing: The Week Ahead's doesn't even include \"Monday\""
        )


def fix_community_time(tbll: list[list[str]]) -> list[list[str]]:
    res = []
    for i in range(1, 5):
        day = tbll[i]
        dayl = [DAYNAMES[i]]
        for j in range(1, len(day)):
            text = day[j]
            if "whole school assembly" in text.lower():
                dayl.append("Whole School Assembly")
            elif (
                "tutor group check-in" in text.lower()
                or "follow up day" in text.lower()
            ):
                dayl.append("Tutor Time")
            else:
                dayl.append(text.strip())
        res.append(dayl)
    return res


def main() -> None:
    config = ConfigParser()
    config.read("config.ini")
    app = msal.PublicClientApplication(
        config["credentials"]["client_id"],
        authority=config["credentials"]["authority"],
    )
    token = acquire_token(app, config)
    download_share_url(
        token,
        config["the_week_ahead"]["file_url"],
        config["the_week_ahead"]["local_filename"],
    )
    data = {
        "community_time_days": fix_community_time(
            extract_community_time_from_presentation(config)
        ),
        "aods": extract_aod_from_presentation(config),
    }
    today = datetime.datetime.today().strftime("%Y%m%d")
    with open(os.path.join("build/", today + "-data.json"), "w") as fd:
        json.dump(data, fd)


if __name__ == "__main__":
    main()
