from __future__ import annotations

from pptx import Presentation # type: ignore
from pprint import pprint

DAYNAMES = ["Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"]

prs = Presentation("the_week_ahead.pptx")

def extract_community_time_from_presentation(prs: Presentation) -> list[list[str]]:
    slide = prs.slides[2]
    for shape in slide.shapes:
        if not shape.has_table:
            continue
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    if row_count != 5 or col_count != 5:
        print("WARNING: irregular table, may cause problems during parsing")
    tbll = []
    for r in range(row_count):
        row = [""]*col_count
        for c in range(col_count):
            cell_text = ""
            cell = tbl.cell(r,c)
            paragraphs = cell.text_frame.paragraphs 
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    cell_text += run.text
            if not cell_text.strip():
                cell_text = old_cell_text # type: ignore
            row[c] = cell_text
            old_cell_text = cell_text
        tbll.append(row)
    return tbll

def fix_community_time(tbll):
    res = []
    for i in range(1, 5):
        day = tbll[i]
        dayl = [DAYNAMES[i]]
        for j in range(1, len(day)):
            text = day[j]
            if "whole school assembly" in text.lower():
                dayl.append("Whole School Assembly")
            elif "tutor group check-in" in text.lower() or "follow up day" in text.lower():
                dayl.append("Tutor Time")
            else:
                dayl.append(text)
        res.append(dayl)
    return res

pprint(fix_community_time(extract_community_time_from_presentation(prs)))
