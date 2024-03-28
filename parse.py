from __future__ import annotations
from typing import Optional, Any
from pprint import pprint
import os
from configparser import ConfigParser
import json
import tempfile

DAYNAMES = ["Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"]

def download_from_sharepoint(sharepoint_site_url: str, sharing_link_url: str, credentials: tuple[str, str], output_to: Optional[str] = None) -> str:
    # TODO: Is the return type correct?
    from office365.sharepoint.client_context import ClientContext # type: ignore
    from office365.runtime.auth.user_credential import UserCredential # type: ignore
    client = ClientContext(sharepoint_site_url).with_credentials(UserCredential(*credentials))
    if not output_to:
        fd, output_to = tempfile.mkstemp()
        # TODO: Is it safe to just discard fd?
    with open(output_to, "wb") as local_file:
        client.web.get_file_by_guest_url(sharing_link_url).download(local_file).execute_query()
    return output_to

def download_the_week_ahead(config: ConfigParser) -> None: # TODO: What happens when the download fails?
    credentials = (config["credentials"]["username"], config["credentials"]["password"])
    sharepoint_site_url = config["the_week_ahead"]["site_url"]
    sharing_link_url = config["the_week_ahead"]["file_url"]
    output_to = config["the_week_ahead"]["local_filename"]
    download_from_sharepoint(sharepoint_site_url, sharing_link_url, credentials, output_to)

def extract_community_time_from_presentation(config: ConfigParser) -> list[list[str]]:
    from pptx import Presentation  # type: ignore
    prs = Presentation(config["the_week_ahead"]["local_filename"])
    slide = prs.slides[int(config["the_week_ahead"]["community_time_page_number"]) - 1]
    for shape in slide.shapes:
        if not shape.has_table:
            continue
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    if row_count != 5 or col_count != 5:
        print("WARNING: irregular table, may cause problems during parsing")
    tbll = []
    for r in range(row_count):
        row = [""] * col_count
        for c in range(col_count):
            cell_text = ""
            cell = tbl.cell(r, c)
            paragraphs = cell.text_frame.paragraphs
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    cell_text += run.text
            if not cell_text.strip():
                cell_text = old_cell_text  # type: ignore
                # TODO: There's probably a more robust way to detect whether a cell spans multiple columns but I'm lazy
            row[c] = cell_text
            old_cell_text = cell_text
        tbll.append(row)
    return tbll

def fix_community_time(tbll: list[list[str]]) -> list[list[str]]:
    res = []
    for i in range(1, 5):
        day = tbll[i]
        dayl = [DAYNAMES[i]]
        for j in range(1, len(day)):
            text = day[j]
            if "whole school assembly" in text.lower():
                dayl.append("Whole School Assembly")
            elif (
                "tutor group check-in" in text.lower()
                or "follow up day" in text.lower()
            ):
                dayl.append("Tutor Time")
            else:
                dayl.append(text.strip())
        res.append(dayl)
    return res

def main():
    config = ConfigParser()
    config.read("config.ini")
    download_the_week_ahead(config)
    original_community_time_table = extract_community_time_from_presentation(config)
    fixed_community_time = fix_community_time(original_community_time_table)
    print(fixed_community_time)

if __name__ == "__main__":
    main()
