#!/usr/bin/env python3
#
# Weekly script to prepare the YK Pao School Daily Bulletin's week JSON data
# Copyright (C) 2024 Runxi Yu <https://runxiyu.org>
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU Affero General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU Affero General Public License for more details.
#
# You should have received a copy of the GNU Affero General Public License
# along with this program.  If not, see <https://www.gnu.org/licenses/>.
#
# Some rules:
# - Pass localized aware datetime objects around.
#   Minimize the use of date strings and numbers.
#   NEVER used naive datetime objects.
#   Frequently check if the tzinfo is correct or cast the zone.
# - Delete variables that aren't supposed to be used anymore.
# - Functions should be short.
# - Do not pass ConfigParser objects around.
# - Use meaningful variable names.
# - Always write type hints.
# - Use the logger! Try not to print.
#
# TODO: Check The Week Ahead's dates

from __future__ import annotations
from typing import Any, Optional, Iterable, Iterator
from configparser import ConfigParser
from pprint import pprint
import argparse
import logging
import subprocess
import datetime
import zoneinfo
import os
import shutil
import json
import base64
import email
import re

import requests
import msal  # type: ignore
import pptx
import pptx.exc
import pypdf

import menuparser

logger = logging.getLogger(__name__)


class MealTableShapeError(ValueError):
    pass


def zero_list(l: list[Any]) -> list[Any]:
    return [(zero_list(i) if (isinstance(i, list)) else "") for i in l]


def equal_shapes(a: list[Any], b: list[Any]) -> bool:
    return zero_list(a) == zero_list(b)


def generate(
    datetime_target: datetime.datetime,  # expected to be local time
    the_week_ahead_url: str,
    the_week_ahead_community_time_page_number: int,
    the_week_ahead_aod_page_number: int,
    weekly_menu_breakfast_page_number: int,
    weekly_menu_lunch_page_number: int,
    weekly_menu_dinner_page_number: int,
    weekly_menu_query_string: str,
    weekly_menu_sender: str,
    weekly_menu_subject_regex: str,
    weekly_menu_subject_regex_four_groups: tuple[int, int, int, int],
    graph_client_id: str,
    graph_authority: str,
    graph_username: str,
    graph_password: str,
    graph_scopes: list[str],
    calendar_address: str,
    soffice: str,
) -> str:
    if not datetime_target.tzinfo:
        raise TypeError("Naive datetimes are unsupported")
    output_filename = "week-%s.json" % datetime_target.strftime("%Y%m%d")
    logger.info("Output filename: %s" % output_filename)

    token: str = acquire_token(
        graph_client_id, graph_authority, graph_username, graph_password, graph_scopes
    )

    calendar_response = requests.get(
        "https://graph.microsoft.com/v1.0/users/%s/calendar/calendarView"
        % calendar_address,
        headers={"Authorization": "Bearer " + token},
        params={
            "startDateTime": datetime_target.replace(microsecond=0).isoformat(),
            "endDateTime": (datetime_target + datetime.timedelta(days=7))
            .replace(microsecond=0)
            .isoformat(),
        },
        timeout=15,
    )
    if calendar_response.status_code != 200:
        raise ValueError(
            "Calendar response status code is not 200", calendar_response.content
        )
    calendar_object = calendar_response.json()
    # pprint(calendar_object)

    the_week_ahead_filename = "the_week_ahead-%s.pptx" % datetime_target.strftime(
        "%Y%m%d"
    )
    if not os.path.isfile(the_week_ahead_filename):
        logger.info(
            "The Week Ahead doesn't seem to exist at %s, downloading"
            % the_week_ahead_filename
        )
        download_share_url(token, the_week_ahead_url, the_week_ahead_filename)
        logger.info("Downloaded The Week Ahead to %s" % the_week_ahead_filename)
        assert os.path.isfile(the_week_ahead_filename)
    else:
        logger.info("The Week Ahead already exists at %s" % the_week_ahead_filename)

    menu_filename = "menu-%s.xlsx" % datetime_target.strftime("%Y%m%d")
    if not (os.path.isfile(menu_filename)):
        logger.info("Menu not found, downloading")
        download_menu(
            token,
            datetime_target,
            weekly_menu_query_string,
            weekly_menu_sender,
            weekly_menu_subject_regex,
            weekly_menu_subject_regex_four_groups,
            menu_filename,
        )
        assert os.path.isfile(menu_filename)
    else:
        logger.info("All menus already exist")

    logger.info("Beginning to parse The Week Ahead")
    the_week_ahead_presentation = pptx.Presentation(the_week_ahead_filename)
    try:
        community_time = extract_community_time(
            the_week_ahead_presentation,
            the_week_ahead_community_time_page_number,
        )
    except ValueError:
        logger.error(
            "Invalid community time! Opening The Week Ahead for manual intervention."
        )
        del the_week_ahead_presentation
        subprocess.run([soffice, the_week_ahead_filename], check=True)
        the_week_ahead_presentation = pptx.Presentation(the_week_ahead_filename)
        community_time = extract_community_time(
            the_week_ahead_presentation,
            the_week_ahead_community_time_page_number,
        )
    del the_week_ahead_filename

    aods = extract_aods(the_week_ahead_presentation, the_week_ahead_aod_page_number)
    # We're assuming the the AODs don't need manual intervention. I think that's fair.
    del the_week_ahead_presentation
    logger.info("Finished parsing The Week Ahead")

    logger.info("Beginning to extract menus")
    menu = menuparser.extract(
        menu_filename,
    )
    logger.info("Finished extracting menus")

    final_data = {
        "start_date": datetime_target.strftime("%Y-%m-%d"),
        "community_time": community_time,
        "aods": aods,
        "menu": menu,
        "snacks": {},  # TODO
    }

    with open(output_filename, "w", encoding="utf-8") as fd:
        json.dump(final_data, fd, ensure_ascii=False, indent="\t")
    logger.info("Dumped to: %s" % output_filename)
    return output_filename


def main() -> None:
    logging.basicConfig(level=logging.INFO)
    parser = argparse.ArgumentParser(description="Weekly script for the Daily Bulletin")
    parser.add_argument(
        "--date",
        default=None,
        help="the start of the week to generate for, in local time, YYYY-MM-DD; defaults to next Monday",
    )
    parser.add_argument(
        "--config", default="config.ini", help="path to the configuration file"
    )
    args = parser.parse_args()

    if args.date:
        datetime_target_naive = datetime.datetime.strptime(args.date, "%Y-%m-%d")
    else:
        datetime_target_naive = None
    del args.date

    config = ConfigParser()
    config.read(args.config)

    tzinfo = zoneinfo.ZoneInfo(config["general"]["timezone"])
    if datetime_target_naive:
        datetime_target_aware = datetime_target_naive.replace(tzinfo=tzinfo)
    else:
        datetime_current_aware = datetime.datetime.now(tz=tzinfo)
        datetime_target_aware = datetime_current_aware + datetime.timedelta(
            days=((-datetime_current_aware.weekday()) % 7)
        )
        del datetime_current_aware
    del datetime_target_naive
    logger.info("Generating for %s" % datetime_target_aware.strftime("%Y-%m-%d %Z"))

    build_path = config["general"]["build_path"]
    # TODO: check if the build path exists and create it if it doesn't
    os.chdir(build_path)

    the_week_ahead_url = config["the_week_ahead"]["file_url"]
    the_week_ahead_community_time_page_number = int(
        config["the_week_ahead"]["community_time_page_number"]
    )
    the_week_ahead_aod_page_number = int(config["the_week_ahead"]["aod_page_number"])

    weekly_menu_breakfast_page_number = int(
        config["weekly_menu"]["breakfast_page_number"]
    )
    weekly_menu_lunch_page_number = int(config["weekly_menu"]["lunch_page_number"])
    weekly_menu_dinner_page_number = int(config["weekly_menu"]["dinner_page_number"])
    weekly_menu_query_string = config["weekly_menu"]["query_string"]
    weekly_menu_sender = config["weekly_menu"]["sender"]
    weekly_menu_subject_regex = config["weekly_menu"]["subject_regex"]
    weekly_menu_subject_regex_four_groups_raw = config["weekly_menu"][
        "subject_regex_four_groups"
    ].split(" ")
    weekly_menu_subject_regex_four_groups = tuple(
        [int(z) for z in weekly_menu_subject_regex_four_groups_raw]
    )
    assert len(weekly_menu_subject_regex_four_groups) == 4
    del weekly_menu_subject_regex_four_groups_raw
    # weekly_menu_dessert_page_number = config["weekly_menu"]["dessert_page_number"]

    graph_client_id = config["credentials"]["client_id"]
    graph_authority = config["credentials"]["authority"]
    graph_username = config["credentials"]["username"]
    graph_password = config["credentials"]["password"]
    graph_scopes = config["credentials"]["scope"].split(" ")

    calendar_address = config["calendar"]["address"]

    soffice = config["general"]["soffice"]

    # TODO: make a function that checks the configuration

    generate(
        datetime_target=datetime_target_aware,
        the_week_ahead_url=the_week_ahead_url,
        the_week_ahead_community_time_page_number=the_week_ahead_community_time_page_number,
        the_week_ahead_aod_page_number=the_week_ahead_aod_page_number,
        weekly_menu_breakfast_page_number=weekly_menu_breakfast_page_number,
        weekly_menu_lunch_page_number=weekly_menu_lunch_page_number,
        weekly_menu_dinner_page_number=weekly_menu_dinner_page_number,
        weekly_menu_query_string=weekly_menu_query_string,
        weekly_menu_sender=weekly_menu_sender,
        weekly_menu_subject_regex=weekly_menu_subject_regex,
        weekly_menu_subject_regex_four_groups=weekly_menu_subject_regex_four_groups,
        graph_client_id=graph_client_id,
        graph_authority=graph_authority,
        graph_username=graph_username,
        graph_password=graph_password,
        graph_scopes=graph_scopes,
        calendar_address=calendar_address,
        soffice=soffice,
    )
    # NOTE: generate() can get the timezone from datetime_target_aware
    # It returns the generated filename.


def encode_sharing_url(url: str) -> str:
    return "u!" + base64.urlsafe_b64encode(url.encode("utf-8")).decode("ascii").rstrip(
        "="
    )


def download_share_url(
    token: str, url: str, local_filename: str, chunk_size: int = 65536
) -> None:

    download_direct_url = requests.get(
        "https://graph.microsoft.com/v1.0/shares/%s/driveItem"
        % encode_sharing_url(url),
        headers={"Authorization": "Bearer " + token},
        timeout=20,
    ).json()["@microsoft.graph.downloadUrl"]

    with requests.get(
        download_direct_url,
        headers={
            "Authorization": "Bearer %s" % token,
            "Accept-Encoding": "identity",
        },
        stream=True,
        timeout=20,
    ) as r:
        with open(local_filename, "wb") as fd:
            shutil.copyfileobj(r.raw, fd)
            fd.flush()


def acquire_token(
    graph_client_id: str,
    graph_authority: str,
    graph_username: str,
    graph_password: str,
    graph_scopes: list[str],
) -> str:
    app = msal.PublicClientApplication(
        graph_client_id,
        authority=graph_authority,
    )
    result = app.acquire_token_by_username_password(
        graph_username, graph_password, scopes=graph_scopes
    )

    if "access_token" in result:
        assert isinstance(result["access_token"], str)
        return result["access_token"]
    raise ValueError("Authentication error in password login")


def search_mail(token: str, query_string: str) -> list[dict[str, Any]]:
    hits = requests.post(
        "https://graph.microsoft.com/v1.0/search/query",
        headers={"Authorization": "Bearer " + token},
        json={
            "requests": [
                {
                    "entityTypes": ["message"],
                    "query": {"queryString": query_string},
                    "from": 0,
                    "size": 15,
                    "enableTopResults": True,
                }
            ]
        },
        timeout=20,
    ).json()["value"][0]["hitsContainers"][0]["hits"]
    assert isinstance(hits, list)
    assert isinstance(hits[0], dict)
    return hits


def extract_aods(
    prs: pptx.presentation.Presentation, aod_page_number: int
) -> list[str]:
    slide = prs.slides[aod_page_number]
    aods = ["", "", "", ""]
    for shape in slide.shapes:
        if hasattr(shape, "text") and "Monday: " in shape.text:
            slist = shape.text.split("\n")
            for s in slist:
                try:
                    day, aod = s.split(": ", 1)
                except ValueError:
                    pass
                day = day.lower()
                if day == "monday":
                    aods[0] = aod
                elif day == "tuesday":
                    aods[1] = aod
                elif day == "wednesday":
                    aods[2] = aod
                elif day == "thursday":
                    aods[3] = aod
            if not all(aods):
                raise ValueError(
                    "AOD parsing: The Week Ahead doesn't include all AOD days, or the formatting is borked"
                )
            return aods
    raise ValueError("AOD parsing: The Week Ahead's doesn't even include \"Monday\"")
    # TODO: this is one of those places where Monday is *expected* to be the first day.
    # TODO: revamp this. this is ugly!


def extract_community_time(
    prs: pptx.presentation.Presentation, community_time_page_number: int
) -> list[list[str]]:

    slide = prs.slides[community_time_page_number]
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        break
    else:
        raise ValueError("No shapes")
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    if col_count not in [4, 5]:
        raise ValueError(
            "Community time parsing: The Week Ahead community time table does not have 4 or 5 columns"
        )
    if col_count == 4:
        logger.warning(
            "Community time warning: only four columns found, assuming that Y12 has graduated"
        )

    res = [["" for c in range(col_count)] for r in range(row_count)]

    for r in range(row_count):
        for c in range(col_count):
            cell = tbl.cell(r, c)
            if not cell.is_spanned:
                t = ""
                for p in cell.text_frame.paragraphs:
                    for pr in p.runs:
                        t += pr.text
                t = t.strip()
                if "whole school assembly" in t.lower():
                    t = "Whole School Assembly"
                elif (
                    "tutor group check-in" in t.lower()
                    or "follow up day" in t.lower()
                    or "open session for tutor and tutee" in t.lower()
                ):
                    t = "Tutor Time"
                res[r][c] = t
                if cell.is_merge_origin:
                    for sh in range(cell.span_height):
                        for sw in range(cell.span_width):
                            res[r + sh][c + sw] = t

    return [x[1:] for x in res[1:]]


def filter_mail_results_by_sender(
    original: Iterable[dict[str, Any]], sender: str
) -> Iterator[dict[str, Any]]:
    for hit in original:
        if (
            hit["resource"]["sender"]["emailAddress"]["address"].lower()
            == sender.lower()
        ):
            yield hit


# TODO: Potentially replace this with a pattern-match based on strptime().
def filter_mail_results_by_subject_regex_groups(
    original: Iterable[dict[str, Any]],
    subject_regex: str,
    subject_regex_groups: Iterable[int],
) -> Iterator[tuple[dict[str, Any], list[str]]]:
    for hit in original:
        logging.debug("Trying %s" % hit["resource"]["subject"])
        matched = re.compile(subject_regex).match(hit["resource"]["subject"])
        if matched:
            yield (hit, [matched.group(group) for group in subject_regex_groups])


def download_menu(
    token: str,
    datetime_target: datetime.datetime,
    weekly_menu_query_string: str,
    weekly_menu_sender: str,
    weekly_menu_subject_regex: str,
    weekly_menu_subject_regex_four_groups: tuple[int, int, int, int],
    menu_filename: str,
) -> None:
    search_results = search_mail(token, weekly_menu_query_string)

    for hit, matched_groups in filter_mail_results_by_subject_regex_groups(
        filter_mail_results_by_sender(search_results, weekly_menu_sender),
        weekly_menu_subject_regex,
        weekly_menu_subject_regex_four_groups,
    ):
        try:
            subject_1st_month = datetime.datetime.strptime(
                matched_groups[0], "%b"  # issues here are probably locales
            ).month
            subject_1st_day = int(matched_groups[1])
        except ValueError:
            raise ValueError(hit["resource"]["subject"], matched_groups[0])
        if (
            subject_1st_month == datetime_target.month
            and subject_1st_day == datetime_target.day
        ):
            break
    else:
        raise ValueError("No SJ-menu email found")

    with requests.get(
        "https://graph.microsoft.com/v1.0/me/messages/%s/$value" % hit["hitId"],
        headers={
            "Authorization": "Bearer %s" % token,
            "Accept-Encoding": "identity",
        },
        stream=True,
        timeout=20,
    ) as r:
        msg = email.message_from_bytes(r.content)

    for part in msg.walk():
        if part.get_content_type() in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ]:
            payload = part.get_payload(decode=True)
            payload_filename_encoded = part.get_filename()
            if not payload_filename_encoded:
                raise ValueError("xlsx does not have a filename")
            payload_filename_mix = email.header.decode_header(payload_filename_encoded)
            assert len(payload_filename_mix) == 1
            payload_filename_encoded, payload_filename_encoding = payload_filename_mix[
                0
            ]
            if payload_filename_encoding is None:
                assert isinstance(payload_filename_encoded, str)
                filename = payload_filename_encoded
            elif isinstance(payload_filename_encoded, bytes):
                filename = payload_filename_encoded.decode(payload_filename_encoding)
            else:
                raise TypeError("What?")

            pb = bytes(payload)

            with open(menu_filename, "wb") as w:
                w.write(pb)
    else:
        raise ValueError("No proper attachment found in email")


if __name__ == "__main__":
    main()
